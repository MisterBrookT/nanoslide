#!/usr/bin/env python3
"""nanoslide CLI - Transform PDFs into engaging slide presentations and videos."""

import json
from pathlib import Path
from typing import Annotated

import typer

from nanoslide.prompts import get_plan_prompt, get_slide_prompt
from nanoslide.utils.google_caller import GoogleCaller
from nanoslide.utils.io import extract_from_markdown

app = typer.Typer(
    name="nanoslide",
    help="Transform PDFs into engaging slide presentations and videos.",
    add_completion=False,
)


def get_output_dir(pdf_path: Path, output_dir: Path) -> Path:
    """Get the output directory for a given PDF file.

    Creates directory structure: output_dir/<pdf_stem>/
    """
    stem = pdf_path.stem
    result_dir = output_dir / stem
    result_dir.mkdir(parents=True, exist_ok=True)
    return result_dir


def _create_pptx(slides_dir: Path, pptx_path: Path) -> None:
    """Create PowerPoint presentation from slide images."""
    from pptx import Presentation
    from pptx.util import Inches

    # Find all slide images
    slide_files = sorted(
        slides_dir.glob("slide_*.png"),
        key=lambda x: int(x.stem.split("_p")[-1])
        if "_p" in x.stem and x.stem.split("_p")[-1].isdigit()
        else 0,
    )

    if not slide_files:
        typer.echo(f"  ⚠️  No slide images found in {slides_dir}")
        return

    typer.echo(f"📑 Creating PowerPoint from {len(slide_files)} slides...")

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 width
    prs.slide_height = Inches(7.5)  # 16:9 height

    # Add each slide image
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_file in slide_files:
        typer.echo(f"  Adding: {slide_file.name}")
        slide = prs.slides.add_slide(blank_layout)

        # Add image to fill the entire slide
        slide.shapes.add_picture(
            str(slide_file),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Save presentation
    prs.save(str(pptx_path))
    typer.echo(f"✅ PowerPoint saved to: {pptx_path}")


@app.command()
def plan(
    pdf: Annotated[Path, typer.Argument(help="Path to the source PDF file")],
    output_dir: Annotated[
        Path, typer.Option("--output", "-o", help="Output directory")
    ] = Path("outputs"),
    prompt: Annotated[
        str | None, typer.Option("--prompt", "-p", help="Style prompt for slides")
    ] = None,
    exist: Annotated[
        bool, typer.Option("--exist", "-e", help="Skip if output already exists")
    ] = False,
) -> Path:
    """Generate a slide plan from a PDF document.

    Analyzes the PDF and creates a detailed JSON plan for each slide.
    Output: outputs/<pdf_name>/plan.json
    """
    if not pdf.exists():
        typer.echo(f"Error: PDF file not found: {pdf}", err=True)
        raise typer.Exit(1)

    # Setup output directory
    result_dir = get_output_dir(pdf, output_dir)
    plan_file = result_dir / "plan.json"

    # Check if already exists
    if exist and plan_file.exists():
        typer.echo(f"⏭️  Skipping plan (already exists): {plan_file}")
        return plan_file

    typer.echo(f"📄 Reading PDF: {pdf}")

    # Default style prompt
    user_prompt = prompt or "蜡笔小新动漫手绘风格，信息密度高，动漫人物形象明显，用中文"

    typer.echo("🧠 Generating slide plan...")
    typer.echo(f"   Style: {user_prompt}")

    # Generate plan using Gemini
    plan_prompt = get_plan_prompt(user_prompt=user_prompt)
    response = GoogleCaller.understand_file(
        model="gemini-3-pro-preview",
        prompt=plan_prompt,
        file_path=pdf,
    )

    json.dump(
        extract_from_markdown(response.text),
        plan_file.open("w"),
        indent=2,
        ensure_ascii=False,
    )
    typer.echo(f"✅ Plan saved to: {plan_file}")
    return plan_file


@app.command()
def gen_slide(
    pdf: Annotated[Path, typer.Argument(help="Path to the source PDF file")],
    output_dir: Annotated[
        Path, typer.Option("--output", "-o", help="Output directory")
    ] = Path("outputs"),
    plan_file: Annotated[
        Path | None,
        typer.Option(
            "--plan", help="Path to existing plan.json (auto-detected if not provided)"
        ),
    ] = None,
    exist: Annotated[
        bool, typer.Option("--exist", "-e", help="Skip slides that already exist")
    ] = False,
) -> Path:
    """Generate slide images from a plan and create PowerPoint.

    Reads the plan.json and generates slide images, then fuses them into a PPT.
    Output: outputs/<pdf_name>/slide_pieces/*.png, outputs/<pdf_name>/presentation.pptx
    """
    result_dir = get_output_dir(pdf, output_dir)

    # Find plan file
    if plan_file is None:
        plan_file = result_dir / "plan.json"

    typer.echo(f"📋 Reading plan: {plan_file}")
    plan_content = json.loads(plan_file.read_text())

    # Create slides directory
    slides_dir = result_dir / "slide_pieces"
    slides_dir.mkdir(parents=True, exist_ok=True)

    typer.echo("🎨 Generating slides...")

    # Parse plan and generate slides - only process sX keys
    slide_keys = [k for k in plan_content.keys() if k.startswith("s")]
    sorted_slide_keys = sorted(
        slide_keys,
        key=lambda x: int(x[1:]) if len(x) > 1 and x[1:].isdigit() else 0,
    )

    previous_slide_path = None
    for slide_key in sorted_slide_keys:
        # Extract slide number from "sX" format
        slide_num = (
            int(slide_key[1:]) if len(slide_key) > 1 and slide_key[1:].isdigit() else 0
        )
        slide_path = slides_dir / f"slide_p{slide_num}.png"

        # Check if already exists
        if exist and slide_path.exists():
            typer.echo(f"  ⏭️  Skipping slide {slide_key} (already exists)")
            previous_slide_path = slide_path
            continue

        slide_content = plan_content[slide_key]
        typer.echo(f"  Generating slide {slide_key}...")

        # Check if previous slide exists for style reference
        has_reference = previous_slide_path is not None and previous_slide_path.exists()
        if has_reference:
            typer.echo(
                f"    Using previous slide as style reference: {previous_slide_path.name}"
            )

        slide_prompt = get_slide_prompt(
            slide_content=slide_content,
            has_reference=has_reference,
        )
        # typer.echo(f"genering slide {slide_key}  with the prompt: {slide_prompt}")
        response = GoogleCaller.generate_image(
            model="gemini-3-pro-image-preview",
            prompt=slide_prompt,
            reference_image_path=previous_slide_path if has_reference else None,
        )

        response.image.save(slide_path)
        typer.echo(f"    ✅ Saved: {slide_path}")

        # Update previous slide path for next iteration
        previous_slide_path = slide_path

    typer.echo(f"✅ Slides saved to: {slides_dir}")

    # Create PowerPoint presentation
    pptx_path = result_dir / "presentation.pptx"
    _create_pptx(slides_dir, pptx_path)

    return slides_dir


@app.command()
def gen_video(
    pdf: Annotated[Path, typer.Argument(help="Path to the source PDF file")],
    output_dir: Annotated[
        Path, typer.Option("--output", "-o", help="Output directory")
    ] = Path("outputs"),
    plan_file: Annotated[
        Path | None,
        typer.Option(
            "--plan", help="Path to existing plan.json (auto-detected if not provided)"
        ),
    ] = None,
    exist: Annotated[
        bool,
        typer.Option("--exist", "-e", help="Skip video segments that already exist"),
    ] = False,
) -> Path:
    """Generate video by interpolating between consecutive slide images.

    Reads consecutive pairs of slide images and generates interpolation videos.
    Output: outputs/<pdf_name>/video_interp/*.mp4
    """
    result_dir = get_output_dir(pdf, output_dir)

    # Find plan file
    if plan_file is None:
        plan_file = result_dir / "plan.json"

    if not plan_file.exists():
        typer.echo(f"Error: Plan file not found: {plan_file}", err=True)
        raise typer.Exit(1)

    # Check if slides exist
    slides_dir = result_dir / "slide_pieces"
    if not slides_dir.exists():
        typer.echo(f"Error: Slides directory not found: {slides_dir}", err=True)
        typer.echo("Run 'nanoslide gen-slide' first to generate slides.")
        raise typer.Exit(1)

    typer.echo(f"📋 Reading plan: {plan_file}")
    plan_content = json.loads(plan_file.read_text())

    # Create video directory for interpolation
    video_dir = result_dir / "video"
    video_dir.mkdir(parents=True, exist_ok=True)

    typer.echo("🎬 Generating interpolation videos from vX keys...")

    # Parse plan and generate interpolation videos - only process vX keys
    video_keys = [k for k in plan_content.keys() if k.startswith("v")]
    sorted_video_keys = sorted(
        video_keys,
        key=lambda x: int(x[1:]) if len(x) > 1 and x[1:].isdigit() else 0,
    )

    # Generate videos for each vX key
    for video_key in sorted_video_keys:
        # Extract video number from "vX" format
        video_num = (
            int(video_key[1:]) if len(video_key) > 1 and video_key[1:].isdigit() else 0
        )

        # vX corresponds to transition between sX and s(X+1)
        slide_num1 = video_num
        slide_num2 = video_num + 1

        video_path = video_dir / f"segment_{video_num}.mp4"

        # Check if already exists
        if exist and video_path.exists():
            typer.echo(f"  ⏭️  Skipping video {video_key} (already exists)")
            continue

        # Get corresponding slide images (sX and s(X+1))
        slide1_path = slides_dir / f"slide_p{slide_num1}.png"
        slide2_path = slides_dir / f"slide_p{slide_num2}.png"

        if not slide1_path.exists() or not slide2_path.exists():
            typer.echo(
                f"  ⚠️  Warning: Slide images not found: {slide1_path.name} or {slide2_path.name}"
            )
            continue

        typer.echo(
            f"  Generating video {video_key} (transition s{slide_num1} → s{slide_num2})..."
        )

        # Use video description directly from plan_content["vX"]
        video_content = plan_content[video_key]
        # video_prompt = get_video_prompt(slide_description=video_content)

        typer.echo(
            f"  Generating video {video_key} (transition s{slide_num1} → s{slide_num2}) with the prompt: {video_content} and the two images: {slide1_path} and {slide2_path}"
        )

        #  model="veo-3.1-fast-generate-preview",
        # model = "veo-3.1-generate-preview"
        response = GoogleCaller.generate_video_interpolation(
            model="veo-3.1-generate-preview",
            prompt=video_content,
            image1_path=slide1_path,
            image2_path=slide2_path,
        )

        if response.video:
            video_path.write_bytes(response.video)
            typer.echo(f"    ✅ Saved: {video_path}")
        else:
            typer.echo(f"    ⚠️  Warning: No video generated for {video_key}")

    typer.echo(f"✅ Interpolation videos saved to: {video_dir}")
    return video_dir


@app.command()
def pipe(
    pdf: Annotated[Path, typer.Argument(help="Path to the source PDF file")],
    output_dir: Annotated[
        Path, typer.Option("--output", "-o", help="Output directory")
    ] = Path("outputs"),
    prompt: Annotated[
        str | None, typer.Option("--prompt", "-p", help="Style prompt for slides")
    ] = None,
    video: Annotated[
        bool, typer.Option("--video", "-v", help="Also generate video from slides")
    ] = True,
    exist: Annotated[
        bool, typer.Option("--exist", "-e", help="Skip steps if output already exists")
    ] = False,
) -> None:
    """Run the complete pipeline: plan → gen_slide (with PPT) → (gen_video).

    This is a convenience command that runs all steps in sequence.
    Output: outputs/<pdf_name>/
      - plan.json
      - slide_pieces/*.png
      - presentation.pptx
      - video/*.mp4 or video_interp/*.mp4 (if --video is set)
    """
    typer.echo("🚀 Starting nanoslide pipeline...")
    typer.echo("=" * 50)

    # Step 1: Generate plan
    typer.echo("\n📝 STEP 1: Generating plan...")
    plan_file = plan(pdf=pdf, output_dir=output_dir, prompt=prompt, exist=exist)

    # Step 2: Generate slides (includes PPT creation)
    typer.echo("\n🎨 STEP 2: Generating slides & PowerPoint...")
    gen_slide(pdf=pdf, output_dir=output_dir, plan_file=plan_file, exist=exist)

    # Step 3: Generate video (optional)
    if video:
        typer.echo("\n🎬 STEP 3: Generating interpolation videos...")
        gen_video(pdf=pdf, output_dir=output_dir, plan_file=plan_file, exist=exist)

    # Step 4: Fuse slides and video
    typer.echo("\n🔗 STEP 4: Fusing slides and video...")
    fuse(pdf=pdf, output_dir=output_dir, slides=True, video=True)

    typer.echo("\n" + "=" * 50)
    typer.echo("✨ Pipeline complete!")
    result_dir = get_output_dir(pdf, output_dir)
    typer.echo(f"📁 Output directory: {result_dir}")


def _merge_videos(video_files: list[Path], output_path: Path) -> None:
    """Merge multiple video files into a single video.

    Args:
        video_files: List of video file paths in order.
        output_path: Path to save the merged video.
    """
    try:
        # Optional dependency: moviepy for video merging
        from moviepy.editor import VideoFileClip, concatenate_videoclips  # type: ignore
    except ImportError:
        # Fallback to ffmpeg if moviepy is not available
        import subprocess
        import tempfile

        # Create a temporary file list for ffmpeg concat
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
            for video_file in video_files:
                f.write(f"file '{video_file.absolute()}'\n")
            concat_file = Path(f.name)

        try:
            # Use ffmpeg to concatenate videos
            subprocess.run(
                [
                    "ffmpeg",
                    "-f",
                    "concat",
                    "-safe",
                    "0",
                    "-i",
                    str(concat_file),
                    "-c",
                    "copy",
                    str(output_path),
                ],
                check=True,
                capture_output=True,
            )
            typer.echo(f"    ✅ Merged {len(video_files)} videos using ffmpeg")
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            typer.echo(f"  ⚠️  Error merging videos: {e}", err=True)
            typer.echo(
                "  💡 Install moviepy (pip install moviepy) or ffmpeg for video merging",
                err=True,
            )
            raise typer.Exit(1)
        finally:
            concat_file.unlink()
        return

    # Use moviepy to merge videos
    clips = []
    for video_file in video_files:
        typer.echo(f"    Loading: {video_file.name}")
        clip = VideoFileClip(str(video_file))
        clips.append(clip)

    typer.echo(f"    Concatenating {len(clips)} video clips...")
    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(
        str(output_path),
        codec="libx264",
        audio_codec="aac",
        temp_audiofile="temp-audio.m4a",
        remove_temp=True,
    )

    # Clean up
    for clip in clips:
        clip.close()
    final_clip.close()

    typer.echo(f"    ✅ Merged {len(video_files)} videos using moviepy")


@app.command()
def fuse(
    pdf: Annotated[Path, typer.Argument(help="Path to the source PDF file")],
    output_dir: Annotated[
        Path, typer.Option("--output", "-o", help="Output directory")
    ] = Path("outputs"),
    slides: Annotated[
        bool, typer.Option("--slides", "-s", help="Fuse slides into PowerPoint")
    ] = True,
    video: Annotated[
        bool, typer.Option("--video", "-v", help="Fuse video segments into one video")
    ] = True,
) -> None:
    """Fuse slides and/or video segments into final outputs.

    Merges individual slide images into a PowerPoint presentation and/or
    merges video segments into a single complete video.
    Output:
      - outputs/<pdf_name>/presentation.pptx (if --slides)
      - outputs/<pdf_name>/video/fused.mp4 (if --video)
    """
    result_dir = get_output_dir(pdf, output_dir)

    typer.echo("🔗 Fusing outputs...")
    typer.echo("=" * 50)

    # Fuse slides into PowerPoint
    if slides:
        slides_dir = result_dir / "slide_pieces"
        pptx_path = result_dir / "presentation.pptx"

        if not slides_dir.exists():
            typer.echo(f"  ⚠️  Slides directory not found: {slides_dir}")
            typer.echo("  Run 'nanoslide gen-slide' first to generate slides.")
        else:
            typer.echo("\n📑 Fusing slides into PowerPoint...")
            _create_pptx(slides_dir, pptx_path)

    # Fuse video segments into one video
    if video:
        video_dir = result_dir / "video"
        fused_video_path = video_dir / "fused.mp4"

        if not video_dir.exists():
            typer.echo(f"  ⚠️  Video directory not found: {video_dir}")
            typer.echo("  Run 'nanoslide gen-video' first to generate video segments.")
        else:
            # Find all video segment files
            video_files = sorted(
                video_dir.glob("segment_*.mp4"),
                key=lambda x: int(x.stem.split("_")[-1])
                if x.stem.split("_")[-1].isdigit()
                else 0,
            )

            if not video_files:
                typer.echo(f"  ⚠️  No video segments found in {video_dir}")
            else:
                typer.echo(f"\n🎬 Fusing {len(video_files)} video segments...")
                for video_file in video_files:
                    typer.echo(f"  Found: {video_file.name}")

                _merge_videos(video_files, fused_video_path)
                typer.echo(f"✅ Fused video saved to: {fused_video_path}")

    typer.echo("\n" + "=" * 50)
    typer.echo("✨ Fusion complete!")
    typer.echo(f"📁 Output directory: {result_dir}")


if __name__ == "__main__":
    app()
